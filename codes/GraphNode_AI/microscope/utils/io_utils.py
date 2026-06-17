from __future__ import annotations

import json
import os
import time
import re
from pathlib import Path
from typing import Any, Dict, List
from langchain_community.document_loaders import (
    PyPDFLoader,
    TextLoader,
    UnstructuredMarkdownLoader,
    Docx2txtLoader,
)
from langchain_core.documents import Document


class PptxLoader:
    """English documentation."""

    def __init__(self, file_path: str):
        self.file_path = file_path

    def load(self) -> List[Document]:
        from pptx import Presentation  # python-pptx

        prs = Presentation(self.file_path)
        docs = []
        for i, slide in enumerate(prs.slides, start=1):
            parts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            parts.append(text)
                elif shape.has_table:
                    for row in shape.table.rows:
                        row_text = " | ".join(
                            cell.text.strip() for cell in row.cells if cell.text.strip()
                        )
                        if row_text:
                            parts.append(row_text)
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    parts.append(f"[Notes] {notes}")
            if parts:
                docs.append(Document(
                    page_content="\n".join(parts),
                    metadata={"source": self.file_path, "slide": i},
                ))
        return docs or [Document(page_content="", metadata={"source": self.file_path})]


def _load_docx_text(file_path: str) -> str:
    try:
        from docx import Document as DocxDocument
    except ImportError as exc:
        raise ImportError("python-docx is required to load .docx files") from exc

    doc = DocxDocument(file_path)
    parts: list[str] = []

    for paragraph in doc.paragraphs:
        text = paragraph.text.strip()
        if text:
            parts.append(text)

    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))

    return "\n\n".join(parts)

def load_ontology_schema(schema_path: str | None = None) -> Dict:
    """
    Load the ontology schema used for extraction and validation.
    """
    if schema_path is None:
        raise ValueError("schema_path is required")
    path = Path(schema_path)
    if not path.is_absolute():
        path = Path(__file__).resolve().parent.parent / path # microscope/path
    with path.open("r", encoding="utf-8") as f:
        return json.load(f)


def load_document(file_path: str) -> List[Document]:
    """English documentation."""
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")

    loader_mapping = {
        ".pdf": PyPDFLoader,
        ".md": UnstructuredMarkdownLoader,
        ".docx": Docx2txtLoader,
        ".pptx": PptxLoader,
    }
    ext = os.path.splitext(file_path)[1].lower()

    if ext not in loader_mapping and ext not in {".txt", ".docx"}:
        raise ValueError(f"Unsupported file extension: {ext}")

    if ext == ".txt":
        encodings = ["utf-8", "utf-8-sig", "cp949", "gbk", "latin-1"]
        last_error = None
        text = None
        for enc in encodings:
            try:
                with open(file_path, "r", encoding=enc) as f:
                    text = f.read()
                break
            except UnicodeDecodeError as exc:
                last_error = exc
        if text is None:
            raise RuntimeError(f"Error loading {file_path}") from last_error
        raw_docs = [Document(page_content=text, metadata={"source": file_path})]
    elif ext == ".docx":
        text = _load_docx_text(file_path)
        raw_docs = [Document(page_content=text, metadata={"source": file_path})]
    else:
        loader = loader_mapping[ext](file_path)
        raw_docs = loader.load()

    # English comment.
    combined_text = "\n\n".join(doc.page_content for doc in raw_docs)
    return [raw_docs[0].__class__(page_content=combined_text, metadata=raw_docs[0].metadata)]

def extract_json_from_text(text: str) -> Any:
    """
    Extract JSON (object or array) from a string that may contain extra text.
    """
    if not text:
        return None
    # If the model returned a JSON string literal (escaped), unescape it first.
    stripped = text.strip()
    if (
        (stripped.startswith('"') and stripped.endswith('"'))
        or (stripped.startswith("'") and stripped.endswith("'"))
    ):
        try:
            text = json.loads(stripped)
        except json.JSONDecodeError:
            text = stripped.strip('"').strip("'")

    code_block_pattern = r"```(?:json)?\s*([\s\S]*?)```"
    code_match = re.search(code_block_pattern, text)
    if code_match:
        text = code_match.group(1).strip()

    try:
        return json.loads(text)
    except json.JSONDecodeError:
        pass

    # English comment.
    try:
        return json.loads(text.replace('\\"', '"'))
    except json.JSONDecodeError:
        pass

    obj_start = text.find("{")
    arr_start = text.find("[")
    if obj_start == -1 and arr_start == -1:
        return None

    start_idx = obj_start if arr_start == -1 else arr_start
    if obj_start != -1 and arr_start != -1:
        start_idx = min(obj_start, arr_start)

    stack = []
    for i in range(start_idx, len(text)):
        ch = text[i]
        if ch in "{[":
            stack.append(ch)
        elif ch in "}]":
            if not stack:
                continue
            open_ch = stack.pop()
            if not stack:
                snippet = text[start_idx : i + 1]
                try:
                    return json.loads(snippet)
                except json.JSONDecodeError:
                    return None
    return None


def save_extraction_outputs(
    output_dir: str,
    source_path: str,
    raw_llm_outputs: List[str],
    extracted_graphs: List[Dict[str, Any]],
    standardized_graphs: List[Dict[str, Any]],
    name_mapping: Dict[str, str],
    chunk_id_map: Dict[int, str],
) -> None:
    os.makedirs(output_dir, exist_ok=True)

    base = os.path.splitext(os.path.basename(source_path))[0]
    raw_path = os.path.join(output_dir, f"{base}_raw_llm.json")
    extracted_path = os.path.join(output_dir, f"{base}_extracted.json")
    standardized_path = os.path.join(output_dir, f"{base}_standardized.json")
    mapping_path = os.path.join(output_dir, f"{base}_name_mapping.json")
    chunk_map_path = os.path.join(output_dir, f"{base}_chunk_id_map.json")

    with open(raw_path, "w", encoding="utf-8") as f:
        json.dump(raw_llm_outputs, f, indent=2, ensure_ascii=True)

    with open(extracted_path, "w", encoding="utf-8") as f:
        json.dump(extracted_graphs, f, indent=2, ensure_ascii=True)

    with open(standardized_path, "w", encoding="utf-8") as f:
        json.dump(standardized_graphs, f, indent=2, ensure_ascii=True)

    with open(mapping_path, "w", encoding="utf-8") as f:
        json.dump(name_mapping, f, indent=2, ensure_ascii=True)

    with open(chunk_map_path, "w", encoding="utf-8") as f:
        json.dump(chunk_id_map, f, indent=2, ensure_ascii=True)


def save_partial_outputs(
    output_dir: str,
    source_path: str,
    raw_llm_outputs: List[str],
    extracted_graphs: List[Dict[str, Any]],
    chunk_id_map: Dict[int, str],
) -> None:
    os.makedirs(output_dir, exist_ok=True)

    base = os.path.splitext(os.path.basename(source_path))[0]
    raw_path = os.path.join(output_dir, f"{base}_raw_llm.json")
    extracted_path = os.path.join(output_dir, f"{base}_extracted.json")
    chunk_map_path = os.path.join(output_dir, f"{base}_chunk_id_map.json")

    with open(raw_path, "w", encoding="utf-8") as f:
        json.dump(raw_llm_outputs, f, indent=2, ensure_ascii=True)

    with open(extracted_path, "w", encoding="utf-8") as f:
        json.dump(extracted_graphs, f, indent=2, ensure_ascii=True)

    with open(chunk_map_path, "w", encoding="utf-8") as f:
        json.dump(chunk_id_map, f, indent=2, ensure_ascii=True)


def save_service_output(service_name: str, payload: Dict[str, Any]) -> Path:
    out_dir = Path(__file__).resolve().parents[2] / "output_data" / "services"
    out_dir.mkdir(parents=True, exist_ok=True)
    stamp = time.strftime("%Y%m%d_%H%M%S")
    path = out_dir / f"{service_name}_{stamp}.json"
    path.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")
    return path
