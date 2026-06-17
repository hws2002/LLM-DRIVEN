"""Microscope pipeline entry point."""

from __future__ import annotations

import base64
import json
import uuid
from pathlib import Path
from typing import Any, Optional

from langchain_core.documents import Document

from infra.repositories.graph.graphnode_repository import GraphNodeDBHandler
from dto.microscope_dto import ToMicroObjectContext
from microscope.graph_generation.generator import GraphGenerator
from microscope.utils.document_utils import chunk_document
from microscope.utils.io_utils import load_document, load_ontology_schema
from shared.api_provider import ApiProvider
from shared.logger import get_shared_logger
from shared.cost_calculator import save_token_run
import shared.config as cfg

logger = get_shared_logger("microscope.call")

_IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp", ".gif"}

_IMAGE_MIME = {
    "jpg": "image/jpeg", "jpeg": "image/jpeg",
    "png": "image/png", "webp": "image/webp", "gif": "image/gif",
}

_IMAGE_CAPTION_SYSTEM_PROMPT = """\
You are an assistant helping build a knowledge graph from educational materials.
Your task: decide if the image contains meaningful educational content worth extracting.

Respond with exactly: SKIP
if the image is any of the following — even if it contains some text or recognizable shape:
- Background textures or patterns (lined paper, graph paper, fabric, wood, etc.)
- Simple geometric shapes (circles, rectangles, arrows used as design elements)
- Decorative borders, frames, or ornamental graphics
- Slide template elements (number markers like "1." "2." used as bullet/section dividers, divider lines, watermarks)
- Simple doodles, cartoon faces, clip-art icons with no domain content
- Logos or branding images without substantive information
- Blank or near-blank slides

Otherwise, write a brief caption (2-4 sentences) describing the educational content — concepts, data, diagrams, charts, photographs, or text slides with meaningful subject matter. Write in the same language as the image content."""


_VISION_SUPPORTED_MIME = {"image/png", "image/jpeg", "image/webp", "image/gif"}


def _call_vision(img_bytes: bytes, mime: str, api_provider: ApiProvider) -> str:
    """English documentation."""
    b64 = base64.b64encode(img_bytes).decode("utf-8")
    messages = [
        {"role": "system", "content": _IMAGE_CAPTION_SYSTEM_PROMPT},
        {"role": "user", "content": [
            {"type": "text", "text": "Write a brief caption for this image."},
            {"type": "image_url", "image_url": {"url": f"data:{mime};base64,{b64}"}},
        ]},
    ]
    return api_provider.chat_completion_text(messages=messages)


def _caption_image_file(file_path: str, api_provider: ApiProvider) -> Document:
    """English documentation."""
    with open(file_path, "rb") as f:
        img_bytes = f.read()
    ext = Path(file_path).suffix.lower().lstrip(".")
    mime = _IMAGE_MIME.get(ext, "image/png")
    caption = _call_vision(img_bytes, mime, api_provider)
    return Document(page_content=caption, metadata={"source": file_path})


def _caption_docx_embedded_images(
    file_path: str, api_provider: ApiProvider
) -> tuple[str, dict]:
    """English documentation."""
    import concurrent.futures
    from docx import Document as DocxDocument

    doc = DocxDocument(file_path)

    tasks = []
    img_num = 0
    for rel in doc.part.rels.values():
        if "image" not in rel.reltype:
            continue
        img_part = rel.target_part
        mime = (img_part.content_type or "").lower()
        if mime not in _VISION_SUPPORTED_MIME:
            logger.debug("DOCX image %d: skipping unsupported format '%s'", img_num + 1, mime)
            continue
        img_num += 1
        tasks.append((img_num, img_part.blob, mime))

    if not tasks:
        return "", {}

    logger.info("DOCX text text %dtext text started", len(tasks))

    results: dict[int, dict | None] = {}

    def _caption(task):
        num, blob, mime = task
        caption = _call_vision(blob, mime, api_provider)
        if caption.strip().upper() == "SKIP":
            logger.info("  Doc Image %d → SKIP (decorative)", num)
            return num, None
        logger.info("  Doc Image %d completed (%d chars)", num, len(caption))
        return num, {
            "caption": caption,
            "base64": base64.b64encode(blob).decode("utf-8"),
            "mime": mime,
        }

    with concurrent.futures.ThreadPoolExecutor(max_workers=5) as executor:
        for num, data in executor.map(_caption, tasks):
            results[num] = data

    lines = []
    image_map = {}
    for num, _, _ in tasks:
        data = results.get(num)
        if data:
            lines.append(f'<image page="0" index="{num}">\n{data["caption"]}\n</image>')
            image_map[(0, num)] = data

    return "\n\n".join(lines), image_map


def _caption_pptx_embedded_images(
    file_path: str, api_provider: ApiProvider
) -> tuple[str, dict]:
    """English documentation."""
    import concurrent.futures
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(file_path)

    tasks = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        img_num = 0
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img = shape.image
                mime = (img.content_type or "").lower()
                if mime not in _VISION_SUPPORTED_MIME:
                    logger.debug("Slide %d: skipping unsupported image format '%s'", slide_num, mime)
                    continue
                img_num += 1
                tasks.append((slide_num, img_num, img.blob, mime))

    if not tasks:
        return "", {}

    logger.info("PPTX text text %dtext text started", len(tasks))

    results: dict[tuple, dict | None] = {}

    def _caption(task):
        slide_num, img_num, blob, mime = task
        caption = _call_vision(blob, mime, api_provider)
        if caption.strip().upper() == "SKIP":
            logger.info("  Slide %d Image %d → SKIP (decorative)", slide_num, img_num)
            return (slide_num, img_num), None
        logger.info("  Slide %d Image %d completed (%d chars)", slide_num, img_num, len(caption))
        return (slide_num, img_num), {
            "caption": caption,
            "base64": base64.b64encode(blob).decode("utf-8"),
            "mime": mime,
        }

    with concurrent.futures.ThreadPoolExecutor(max_workers=5) as executor:
        for key, data in executor.map(_caption, tasks):
            results[key] = data

    lines = []
    image_map = {}
    for slide_num, img_num, _, _ in tasks:
        data = results.get((slide_num, img_num))
        if data:
            lines.append(f'<image page="{slide_num}" index="{img_num}">\n{data["caption"]}\n</image>')
            image_map[(slide_num, img_num)] = data

    return "\n\n".join(lines), image_map


class DocumentTooLargeError(ValueError):
    """English documentation."""
    error_code = "DOCUMENT_TOO_LARGE"

_SCHEMA_DIR = Path(__file__).resolve().parent / "schema"
_DEFAULT_SCHEMA = "ontology_schema_general.json"


def _resolve_schema_path(schema_name: Optional[str]) -> str:
    if schema_name:
        candidate = _SCHEMA_DIR / f"ontology_schema_{schema_name}.json"
        if candidate.exists():
            return str(candidate)
        logger.warning("Schema '%s' not found, falling back to default", schema_name)
    return str(_SCHEMA_DIR / _DEFAULT_SCHEMA)


_TYPE_MAPPING_PATH = _SCHEMA_DIR / "type_mapping.json"


def _build_type_mapping(user_language: Optional[str]) -> dict:
    """English documentation."""
    if not user_language:
        return {}
    if not _TYPE_MAPPING_PATH.exists():
        logger.warning("type_mapping.json not found — skipping type mapping")
        return {}
    full_mapping = load_ontology_schema(str(_TYPE_MAPPING_PATH))
    mapping = {}
    for en_name, langs in full_mapping.get("node_types", {}).items():
        if user_language in langs:
            mapping[en_name] = langs[user_language]
    for en_name, langs in full_mapping.get("edge_types", {}).items():
        if user_language in langs:
            mapping[en_name] = langs[user_language]
    return mapping


def _apply_type_mapping(results: list, mapping: dict) -> list:
    """English documentation."""
    if not mapping:
        return results
    for batch in results:
        for node in batch.get("nodes", []):
            t = node.get("type")
            if isinstance(t, list):
                node["type"] = [mapping.get(x, x) for x in t]
            elif isinstance(t, str):
                node["type"] = mapping.get(t, t)
        for edge in batch.get("edges", []):
            t = edge.get("type")
            if isinstance(t, str):
                edge["type"] = mapping.get(t, t)
    return results


def call(
    *,
    to_micro: ToMicroObjectContext,
    api_provider: ApiProvider,
    graph_store: GraphNodeDBHandler,
    schema_name: Optional[str] = None,
    schema: Optional[dict] = None,
    user_language: Optional[str] = None,
    block_mode: bool = False,
    block_granularity: str = "medium",
    source_type: str = "chat",
    generate_micro_graphs: bool = False,
    save_dir: Optional[Path] = None,
    skip_store: bool = False,
) -> dict[str, Any]:
    """Run end-to-end microscope pipeline.

    Args:
        block_mode       : If True, run BLOCK-based consciousness-flow pipeline
        block_granularity: "coarse" | "medium" | "fine" (only used in block_mode)
        source_type      : "chat" | "note" (only used in block_mode)
        generate_micro_graphs: If True in block_mode, extract nested micro graphs per block.
    """
    if block_mode:
        return _call_block_mode(
            to_micro=to_micro,
            api_provider=api_provider,
            graph_store=graph_store,
            schema_name=schema_name,
            schema=schema,
            granularity=block_granularity,
            source_type=source_type,
            generate_micro_graphs=generate_micro_graphs,
            save_dir=save_dir,
            skip_store=skip_store,
        )
    file_path = to_micro.file_path
    file_name = to_micro.file_name
    user_id = to_micro.user_id
    group_id = to_micro.group_id

    chunk_size = cfg.MICROSCOPE_CHUNK_SIZE
    chunk_overlap = cfg.MICROSCOPE_CHUNK_OVERLAP
    batch_max_tokens = cfg.MICROSCOPE_BATCH_MAX_TOKENS
    custom_schema = schema is not None
    if not custom_schema:
        schema_path = _resolve_schema_path(schema_name)
        schema = load_ontology_schema(schema_path)
        logger.info("Using schema: %s", schema_path)
    else:
        logger.info("Using custom schema (dict)")

    logger.info("Processing document: %s", file_path)
    graph_generator = GraphGenerator(
        api_provider=api_provider,
        user_id=user_id,
        ontology_schema=schema,
        batch_max_tokens=batch_max_tokens,
    )

    ext = Path(file_path).suffix.lower()
    image_map: dict = {}
    if ext in _IMAGE_EXTENSIONS:
        logger.info("Image file detected — extracting caption via Vision API")
        documents = [_caption_image_file(file_path, api_provider)]
        logger.info("Caption extracted (%d chars)", len(documents[0].page_content))
    else:
        documents = load_document(file_path)
        if ext == ".pptx":
            img_caption_text, image_map = _caption_pptx_embedded_images(file_path, api_provider)
            if img_caption_text:
                logger.info("PPTX text text text text (%d chars, %dtext)", len(img_caption_text), len(image_map))
                documents[0] = Document(
                    page_content=documents[0].page_content + "\n\n" + img_caption_text,
                    metadata=documents[0].metadata,
                )
        elif ext == ".docx":
            img_caption_text, image_map = _caption_docx_embedded_images(file_path, api_provider)
            if img_caption_text:
                logger.info("DOCX text text text text (%d chars, %dtext)", len(img_caption_text), len(image_map))
                documents[0] = Document(
                    page_content=documents[0].page_content + "\n\n" + img_caption_text,
                    metadata=documents[0].metadata,
                )
    logger.info("Loaded %s document(s)", len(documents))

    chunks = chunk_document(documents, chunk_size, chunk_overlap)
    if not chunks:
        raise ValueError("No chunks produced")
    logger.info("Created %s chunks", len(chunks))

    # English comment.
    estimated_batches = max(1, len(chunks) // max(1, batch_max_tokens // chunk_size))
    MAX_BATCHES = 30
    if estimated_batches >= MAX_BATCHES:
        raise DocumentTooLargeError(
            f"text text text: text text text {estimated_batches}text (text {MAX_BATCHES - 1}text). "
            f"text text text."
        )

    extracted_entity_relations, raw_llm_outputs = (
        graph_generator.extract_entity_relation_from_chunks(chunks, True)
    )
    logger.info("Extracted %s batches", len(extracted_entity_relations))

    logger.info("Standardizing entities...")
    existing_nodes = []
    try:
        existing_nodes = graph_store.get_nodes_by_group_id(group_id=group_id, user_id=user_id)
    except Exception as exc:
        logger.warning("Failed to get existing nodes: %s", exc)

    standardized_results, name_mapping = graph_generator.standardize_extracted_graph(
        extracted_entity_relations,
        existing_nodes=existing_nodes,
    )
    logger.info("Standardization complete")

    source_id = str(uuid.uuid4())
    if skip_store:
        logger.info("skip_store=True — DB text text")
        store_stats = {"chunks_stored": 0, "entities_stored": 0, "edges_stored": 0, "entities_without_chunk_id": 0}
    else:
        logger.info("Ingesting into VectorDB/Neo4j...")
        store_stats = graph_store.store_standardized_data(
            standardized_results=standardized_results,
            source_name=file_name,
            source_id=source_id,
            user_id=user_id,
            group_id=group_id,
            chunks=chunks,
        )
        logger.info(
            "Ingestion complete. Chunks: %s, Entities: %s, Edges: %s",
            store_stats["chunks_stored"],
            store_stats["entities_stored"],
            store_stats["edges_stored"],
        )
        if store_stats["entities_without_chunk_id"] > 0:
            logger.warning(
                "%s entities without source_chunk_id",
                store_stats["entities_without_chunk_id"],
            )

    # English comment.
    if user_language and not custom_schema:
        type_mapping = _build_type_mapping(user_language)
        standardized_results = _apply_type_mapping(standardized_results, type_mapping)

    save_token_run(graph_generator.tracker, source_id, service_name="microscope", user_id=user_id)

    all_nodes: list[dict[str, Any]] = []
    all_edges: list[dict[str, Any]] = []
    for batch in standardized_results:
        all_nodes.extend(batch.get("nodes", []))
        all_edges.extend(batch.get("edges", []))

    # English comment.
    if image_map:
        import re
        _TAG = re.compile(r'<image page="(\d+)" index="(\d+)">')
        chunk_to_images: dict[int, list] = {}
        for idx, chunk in enumerate(chunks):
            keys = [
                (int(p), int(i))
                for p, i in _TAG.findall(chunk.page_content)
                if (int(p), int(i)) in image_map
            ]
            if keys:
                chunk_to_images[idx] = keys

        for node in all_nodes:
            chunk_idx = node.get("source_chunk_id")
            if isinstance(chunk_idx, str) and chunk_idx.isdigit():
                chunk_idx = int(chunk_idx)
            if chunk_idx is not None and chunk_idx in chunk_to_images:
                node["images"] = [
                    {
                        "slide": s,
                        "img": i,
                        "caption": image_map[(s, i)]["caption"],
                        "base64": image_map[(s, i)]["base64"],
                        "mime": image_map[(s, i)]["mime"],
                    }
                    for s, i in chunk_to_images[chunk_idx]
                ]

    chunk_id_map = {
        idx: chunk.metadata.get("chunk_id")
        for idx, chunk in enumerate(chunks)
        if chunk.metadata.get("chunk_id")
    }

    token_usage = graph_generator.tracker.to_dict()

    if save_dir:
        _write_json(save_dir / "standardized.json", standardized_results)
        _write_json(save_dir / "token_usage.json", token_usage)
        logger.info("Dev output saved → %s", save_dir)

    return {
        "nodes": all_nodes,
        "edges": all_edges,
        "source_name": file_name,
        "source_id": source_id,
        "source_path": file_path,
        "chunks_count": len(chunks),
        "batches_count": len(extracted_entity_relations),
        "name_mapping": name_mapping,
        "ingest_stats": store_stats,
        "raw_llm_outputs": raw_llm_outputs,
        "extracted_graphs": extracted_entity_relations,
        "standardized_graphs": standardized_results,
        "token_usage": token_usage,
        "chunk_id_map": chunk_id_map,
    }


def _write_text(path: Path, text: str) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(text, encoding="utf-8")


def _write_json(path: Path, data: Any) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(data, ensure_ascii=False, indent=2), encoding="utf-8")


def _save_block_text_outputs(save_dir: Path, blocks: list) -> None:
    blocks_dir = save_dir / "03_block_texts"
    summary_index = []
    for block in blocks:
        block_dir = blocks_dir / block.block_id
        _write_text(block_dir / "raw_text.txt", block.raw_text)
        summary = {
            "block_id": block.block_id,
            "title": block.title,
            "summary": block.summary,
            "key_concepts": block.key_concepts,
            "order_index": block.order_index,
            "source_type": block.source_type,
        }
        if block.turn_range is not None:
            summary["turn_range"] = list(block.turn_range)
        if block.heading_path is not None:
            summary["heading_path"] = block.heading_path
        _write_json(block_dir / "summary.json", summary)
        summary_index.append(summary)
    _write_json(blocks_dir / "blocks_summary.json", summary_index)


def _call_block_mode(
    *,
    to_micro: ToMicroObjectContext,
    api_provider: ApiProvider,
    graph_store: GraphNodeDBHandler,
    schema_name: Optional[str],
    schema: Optional[dict],
    granularity: str,
    source_type: str,
    generate_micro_graphs: bool = False,
    save_dir: Optional[Path] = None,
    skip_store: bool = False,
) -> dict[str, Any]:
    """BLOCK-based consciousness-flow pipeline with optional progressive save."""
    from microscope.block.segmenter import BlockSegmenter
    from microscope.block.orderer import BlockOrderer
    from microscope.block.assembler import BlockAssembler

    file_path = to_micro.file_path
    file_name = to_micro.file_name
    user_id   = to_micro.user_id

    if schema is None:
        schema_path = _resolve_schema_path(schema_name)
        schema = load_ontology_schema(schema_path)

    graph_generator = GraphGenerator(
        api_provider=api_provider,
        user_id=user_id,
        ontology_schema=schema,
        batch_max_tokens=cfg.MICROSCOPE_BATCH_MAX_TOKENS,
    )

    # English comment.
    documents = load_document(file_path)
    ext = Path(file_path).suffix.lower()

    # English comment.
    # English comment.
    # English comment.
    # English comment.
    image_map: dict = {}
    img_caption_text = ""
    if ext == ".pptx":
        img_caption_text, image_map = _caption_pptx_embedded_images(file_path, api_provider)
    elif ext == ".docx":
        img_caption_text, image_map = _caption_docx_embedded_images(file_path, api_provider)

    if img_caption_text:
        from langchain_core.documents import Document as _Doc
        documents[0] = _Doc(
            page_content=documents[0].page_content + "\n\n" + img_caption_text,
            metadata=documents[0].metadata,
        )

    raw_text = "\n".join(doc.page_content for doc in documents)

    # English comment.
    if save_dir and image_map:
        import re as _re
        images_dir = save_dir / "images"
        images_dir.mkdir(parents=True, exist_ok=True)

        # image_map key: (page_num, img_num) → {caption, base64, mime}
        _ext_map = {"image/png": "png", "image/jpeg": "jpg", "image/webp": "webp", "image/gif": "gif"}
        tag_to_path: dict[tuple, str] = {}
        for (page_num, img_num), data in image_map.items():
            img_ext = _ext_map.get(data["mime"], "png")
            img_filename = f"page_{page_num}_img_{img_num}.{img_ext}"
            img_path = images_dir / img_filename
            img_path.write_bytes(base64.b64decode(data["base64"]))
            tag_to_path[(page_num, img_num)] = f"images/{img_filename}"

        def _replace_tag(match: "re.Match") -> str:
            page_num, img_num = int(match.group(1)), int(match.group(2))
            caption = match.group(3).strip()
            rel_path = tag_to_path.get((page_num, img_num), f"page_{page_num}_img_{img_num}")
            return f"{{{rel_path}}} <caption>{caption}</caption>"

        # English comment.
        raw_text = _re.sub(
            r'<image page="(\d+)" index="(\d+)">\n(.*?)\n</image>',
            _replace_tag,
            raw_text,
            flags=_re.DOTALL,
        )
        logger.info("[block_mode] %d text text → %s", len(image_map), images_dir)

    # ── Step 1: Block Segmentation ────────────────────────────────────────────
    logger.info("[block_mode] Segmenting into blocks (granularity=%s, source_type=%s)", granularity, source_type)
    segmenter = BlockSegmenter(api_provider)
    blocks    = segmenter.segment(raw_text, source_type=source_type, granularity=granularity)
    logger.info("[block_mode] %d blocks created", len(blocks))

    if save_dir:
        seg_dir = save_dir / "01_segmentation"
        _write_text(seg_dir / "prompt_system.txt", segmenter.last_system_prompt)
        _write_text(seg_dir / "prompt_user.txt",   segmenter.last_user_prompt)
        _write_text(seg_dir / "raw_response.txt",  segmenter.last_raw)
        logger.info("[block_mode] Segmentation saved → %s", seg_dir)

    # ── Step 2: Block Ordering ────────────────────────────────────────────────
    logger.info("[block_mode] Analyzing block dependencies...")
    orderer = BlockOrderer(api_provider)
    edges, paths, rationale = orderer.order(blocks)
    logger.info("[block_mode] %d block edges, %d paths", len(edges), len(paths))

    if save_dir:
        ord_dir = save_dir / "02_ordering"
        _write_text(ord_dir / "prompt_system.txt", orderer.last_system_prompt)
        _write_text(ord_dir / "prompt_user.txt",   orderer.last_user_prompt)
        _write_text(ord_dir / "raw_response.txt",  orderer.last_raw)
        logger.info("[block_mode] Ordering saved → %s", ord_dir)

    # ── Step 3: Per-block micro extraction ────────────────────────────────────
    if generate_micro_graphs:
        logger.info("[block_mode] Extracting micro graphs per block...")
        assembler = BlockAssembler(
            graph_generator=graph_generator,
            chunk_size=cfg.MICROSCOPE_CHUNK_SIZE,
            chunk_overlap=cfg.MICROSCOPE_CHUNK_OVERLAP,
        )
        micro_dir = (save_dir / "03_micro_graphs") if save_dir else None
        assembler.extract_micro_graphs(blocks, save_dir=micro_dir)
    else:
        logger.info("[block_mode] Skipping per-block micro graph extraction")

    # ── Step 4: Assemble final BlockGraph ─────────────────────────────────────
    block_graph = BlockAssembler.assemble(blocks, edges, paths, source_type, rationale)
    result      = block_graph.to_dict()
    result["source_id"]   = str(uuid.uuid4())
    result["source_name"] = file_name
    result["source_path"] = file_path
    result["token_usage"] = graph_generator.tracker.to_dict()
    result["image_map"]   = image_map  # English comment.

    if save_dir:
        _save_block_text_outputs(save_dir, blocks)
        _write_json(save_dir / "block_graph.json",  result)
        _write_json(save_dir / "token_usage.json",  result["token_usage"])
        logger.info("[block_mode] Final block graph saved → %s", save_dir)

    return result
