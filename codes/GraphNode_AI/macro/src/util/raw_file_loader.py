"""Load PDF, PPTX, DOCX files into InputData format for macro pipeline."""

from __future__ import annotations

import re
from pathlib import Path
from typing import List, Optional

from util.io_schemas import InputData, Section, SourceNode


# ── PDF ──────────────────────────────────────────────────────────────────────

def load_pdf(path: Path) -> InputData:
    """English documentation."""
    try:
        import fitz  # PyMuPDF
    except ImportError as e:
        raise ImportError("pymupdf is required: pip install pymupdf") from e

    path = Path(path)
    doc = fitz.open(str(path))

    sections: List[Section] = []
    for page_idx in range(len(doc)):
        page = doc[page_idx]
        text = page.get_text("text").strip()
        if not text or len(text) < 10:
            continue
        sections.append(Section(
            id=f"pdf_{path.stem}_p{page_idx}",
            content=text,
            role=None,
            section_title=f"Page {page_idx + 1}",
        ))

    doc.close()

    if not sections:
        return InputData(source_nodes=[])

    safe_name = re.sub(r"[^\w\-]", "_", path.stem)[:40]
    node = SourceNode(
        id=f"pdf_{safe_name}",
        title=path.stem,
        sections=sections,
        source_type="pdf",
        create_time=None,
        update_time=None,
    )
    return InputData(source_nodes=[node])


# ── PPTX ─────────────────────────────────────────────────────────────────────

def load_pptx(path: Path) -> InputData:
    """English documentation."""
    try:
        from pptx import Presentation
    except ImportError as e:
        raise ImportError("python-pptx is required: pip install python-pptx") from e

    path = Path(path)
    prs = Presentation(str(path))

    sections: List[Section] = []
    for slide_idx, slide in enumerate(prs.slides):
        parts: List[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    parts.append(text)

        content = "\n".join(parts).strip()
        if not content or len(content) < 10:
            continue

        sections.append(Section(
            id=f"pptx_{path.stem}_s{slide_idx}",
            content=content,
            role=None,
            section_title=f"Slide {slide_idx + 1}",
        ))

    if not sections:
        return InputData(source_nodes=[])

    safe_name = re.sub(r"[^\w\-]", "_", path.stem)[:40]
    node = SourceNode(
        id=f"pptx_{safe_name}",
        title=path.stem,
        sections=sections,
        source_type="pptx",
        create_time=None,
        update_time=None,
    )
    return InputData(source_nodes=[node])


# ── DOCX ─────────────────────────────────────────────────────────────────────

def load_docx(path: Path) -> InputData:
    """English documentation."""
    try:
        from docx import Document
    except ImportError as e:
        raise ImportError("python-docx is required: pip install python-docx") from e

    path = Path(path)
    doc = Document(str(path))

    sections: List[Section] = []
    current_heading: Optional[str] = None
    current_lines: List[str] = []
    sec_idx = 0

    def _flush(heading: Optional[str], lines: List[str], idx: int) -> Optional[Section]:
        content = "\n".join(lines).strip()
        if not content or len(content) < 10:
            return None
        return Section(
            id=f"docx_{path.stem}_sec{idx}",
            content=content,
            role=None,
            section_title=heading,
        )

    for para in doc.paragraphs:
        style = para.style.name or ""
        text = para.text.strip()
        if not text:
            continue

        if style.startswith("Heading"):
            # English comment.
            sec = _flush(current_heading, current_lines, sec_idx)
            if sec:
                sections.append(sec)
                sec_idx += 1
            current_heading = text
            current_lines = []
        else:
            current_lines.append(text)

    # English comment.
    sec = _flush(current_heading, current_lines, sec_idx)
    if sec:
        sections.append(sec)

    if not sections:
        return InputData(source_nodes=[])

    safe_name = re.sub(r"[^\w\-]", "_", path.stem)[:40]
    node = SourceNode(
        id=f"docx_{safe_name}",
        title=path.stem,
        sections=sections,
        source_type="docx",
        create_time=None,
        update_time=None,
    )
    return InputData(source_nodes=[node])
